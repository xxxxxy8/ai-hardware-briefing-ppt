import os
import json
import re
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ==========================================
# 全局配置参数
# ==========================================
MY_DEEPSEEK_KEY = "sk-d84519a027e64127a169f59827ede10b"   # 密钥
API_URL = "https://api.deepseek.com/v1" 
MODEL_TYPE = "deepseek-chat"           
SRC_FILE = "paper.txt"                  # 原始论文输入
OUT_PPT_NAME = "AI_Hardware_Briefing.pptx" # 导出的PPT文件名


# ==========================================
# 1. 论文文本本地处理模块,减少token消耗
# ==========================================

def load_paper_data(path):
    """
    读取txt文本内容
    """
    if not os.path.exists(path):
        raise FileNotFoundError(f"未找到输入文件: {path}")
    with open(path, "r", encoding="utf-8") as f:
        all_lines = f.readlines()
    return all_lines


def format_and_clean(lines_list):
    """
    去掉空行和两边的空格
    """
    clean_arr = []
    for l in lines_list:
        tmp = l.strip()
        if tmp: 
            clean_arr.append(tmp)
    return clean_arr


def locate_tag(lines, tag_str):
    """
    定位关键词所在的行号
    """
    for idx, l in enumerate(lines):
        if tag_str in l:
            return idx
    return -1


def get_block_context(lines, keyword, pre_len=5, post_len=25):
    """
    截取关键词前后的上下文
    """
    pos = locate_tag(lines, keyword)
    if pos == -1:
        return f"[Warning: 没找到章节标签 -> {keyword}]\n"
    
    low = max(0, pos - pre_len)
    high = min(len(lines), pos + post_len)
    
    return "\n".join(lines[low:high])


def do_text_extract():
    """
    执行本地文本流提取流程
    """
    print(">>> 正在加载原始文本...")
    lines_raw = load_paper_data(SRC_FILE)
    lines_cleaned = format_and_clean(lines_raw)
    
    print(">>> 按照3页大纲要求精准匹配核心模块...")
    buffer = []
    
    # Slide 1: 1000x Vision 概述与核心痛点
    buffer.append("\n=== SLIDE 1 CATEGORY: 1000x Vision Overview ===")
    buffer.append(get_block_context(lines_cleaned, "1000"))
    buffer.append(get_block_context(lines_cleaned, "Unsustainable"))
    
    # Slide 2: 三层抽象架构 (Hardware, Algorithm, Application)
    buffer.append("\n=== SLIDE 2 CATEGORY: Three Abstraction Layers ===")
    buffer.append(get_block_context(lines_cleaned, "Hardware Layer"))
    buffer.append(get_block_context(lines_cleaned, "Algorithm Layer"))
    buffer.append(get_block_context(lines_cleaned, "Application Layer"))
    
    # Slide 3: 前瞻趋势与未来挑战 (Near-Term, Long-Term, Obstacles)
    buffer.append("\n=== SLIDE 3 CATEGORY: Future Trends and Obstacles ===")
    buffer.append(get_block_context(lines_cleaned, "Near-Term"))
    buffer.append(get_block_context(lines_cleaned, "Long-Term"))
    buffer.append(get_block_context(lines_cleaned, "Obstacles"))
    
    return "\n".join(buffer)


# ==========================================
# 2. DeepSeek API 交互与大模型处理
# ==========================================

def request_deepseek_json(context_text):
    """
    调用大模型获取结构化JSON数据，加入严格稳定容错脚手架
    """
    api_client = OpenAI(api_key=MY_DEEPSEEK_KEY, base_url=API_URL)
    
    sys_instruction = (
        "You are a professional PPT structure generator. "
        "Analyze the input text and extract key points into a precise PPT outline. "
        "Strictly output a raw JSON array containing EXACTLY 3 slide objects. "
        "No explanations, no markdown wrappers like ```json.\n\n"
        
        "CRITICAL FORMAT STRUCTURING:\n"
        "Your final output MUST follow this exact schema structure:\n"
        "[\n"
        "  {\n"
        "    \"title\": \"1000X VISION & ECO-SYSTEM\",\n"
        "    \"points\": [\"point 1\", \"point 2\", \"point 3\", \"point 4\"]\n"
        "  },\n"
        "  {\n"
        "    \"title\": \"THREE ABSTRACTION LAYERS\",\n"
        "    \"points\": [\"Hardware Layer: brief summary\", \"Algorithm Layer: brief summary\", \"Application Layer: brief summary\"]\n"
        "  },\n"
        "  {\n"
        "    \"title\": \"FUTURE TRENDS & OBSTACLES\",\n"
        "    \"points\": [\"Near-Term: brief trend\", \"Long-Term: brief trend\", \"Main Obstacle: brief risk\", \"Solution: core strategy\"]\n"
        "  }\n"
        "]\n\n"
        "CRITICAL RULES:\n"
        "- Slide 2 MUST have EXACTLY 3 bullet points (Hardware, Algorithm, Application). Do not add any extra lines.\n"
        "- Slide 3 MUST have EXACTLY 4 bullet points representing Near-term, Long-term, Obstacles, and Solutions.\n"
        "- Every point text must be highly technical, clear, and under 15 words. Never leave 'points' array empty."
    )
    
    user_input = f"Input Text:\n{context_text}"
    
    print(">>> 正在请求 DeepSeek 进行严格3页容错提炼...")
    api_resp = api_client.chat.completions.create(
        model=MODEL_TYPE,
        messages=[
            {"role": "system", "content": sys_instruction},
            {"role": "user", "content": user_input}
        ],
        temperature=0.3 # 恢复到0.3提供必要流利度，依靠结构脚手架锁死格式
    )
    
    res_text = api_resp.choices[0].message.content.strip()
    
    if res_text.startswith("```"):
        res_text = re.sub(r"^```[a-zA-Z]*\n|```$", "", res_text, flags=re.MULTILINE).strip()
        
    return json.loads(res_text)


# ==========================================
# 3. PPTX 自动化渲染与排版引擎
# ==========================================

def build_presentation(json_data, save_path):
    """
    动态渲染完美不溢出、不留白的3页PPT大厂视觉风格
    """
    print(">>> 开始利用 python-pptx 渲染自适应幻灯片...")
    doc = Presentation()
    doc.slide_width = Inches(13.333)
    doc.slide_height = Inches(7.5)
    layout_mode = doc.slide_layouts[6] 
    
    RGB_BG = RGBColor(15, 23, 42)       
    RGB_TITLE = RGBColor(255, 255, 255) 
    RGB_BODY = RGBColor(148, 163, 184)  
    RGB_HIGHLIGHT = RGBColor(56, 189, 248) 
    
    final_slides = json_data[:3]
    
    for i, item in enumerate(final_slides):
        curr_slide = doc.slides.add_slide(layout_mode)
        
        # 1. 全屏背景
        bg_rect = curr_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, doc.slide_width, doc.slide_height)
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = RGB_BG
        bg_rect.line.fill.background()
        
        # 2. 标题
        t_box = curr_slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.5), Inches(1.0))
        t_frame = t_box.text_frame
        t_frame.word_wrap = True
        t_frame.margin_left = t_frame.margin_top = t_frame.margin_right = t_frame.margin_bottom = 0
        
        p_title = t_frame.paragraphs[0]
        p_title.text = item.get("title", f"SECTION {i + 1:02d}").upper()
        p_title.font.size = Pt(38) 
        p_title.font.bold = True
        p_title.font.color.rgb = RGB_TITLE
        p_title.font.name = "Arial"
        
        # 3. 正文框
        b_box = curr_slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.5), Inches(4.5))
        b_frame = b_box.text_frame
        b_frame.word_wrap = True
        b_frame.margin_left = b_frame.margin_top = b_frame.margin_right = b_frame.margin_bottom = 0
        
        points_list = item.get("points", [])
        for p_idx, text_str in enumerate(points_list):
            p_body = b_frame.paragraphs[0] if p_idx == 0 else b_frame.add_paragraph()
            p_body.text = f"▪  {text_str}"
            p_body.font.size = Pt(24) 
            p_body.font.color.rgb = RGB_BODY
            p_body.font.name = "Calibri"
            
            # 页数项数间距控制
            if len(points_list) == 3:
                p_body.space_after = Pt(36)
            else:
                p_body.space_after = Pt(22)
            
            # 第一行默认赋予高亮
            if p_idx == 0:
                p_body.font.color.rgb = RGB_HIGHLIGHT
                p_body.font.bold = True

    doc.save(save_path)
    print(f"🎉 任务搞定！满血版 3 页不留白 PPT 已成功输出至: {save_path}")


# ==========================================
# 主程序入口
# ==========================================
if __name__ == "__main__":
    try:
        extracted_text = do_text_extract()
        outline_json = request_deepseek_json(extracted_text)
        build_presentation(outline_json, OUT_PPT_NAME)
    except Exception as err:
        print(f"\n[Error] 脚本运行期间发生异常: {err}")
